from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import PP_MEDIA_TYPE
from os import listdir
from os.path import isfile, join


prs = Presentation()

# insert .mp3 and .png files in src folder
src = '.\\src\\'
fileNames = [f for f in listdir(src) if isfile(join(src, f))]

# print files sorted
# print(fileNames)

size = len(fileNames)
for i in range(0, size, 2):

    # chose a slide layout to use
    slide_register = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_register)

    # add audio to the slide
    audio_shape = slide.shapes.add_movie(
        src + fileNames[i],
        left=Inches(4.17),
        top=Inches(6.74),
        width=Inches(1.67),
        height=Inches(0.76),
        poster_frame_image=None,
        mime_type='audio/mp3')

    # audio_shape.media_type = PP_MEDIA_TYPE.AUDIO

    img_path = src + fileNames[i + 1]

    # add_picture(image_file, left, top, width, height)
    pic = slide.shapes.add_picture(img_path, Inches(
        0.5), Inches(1.75), width=Inches(9), height=Inches(5))

# save file as example.pptx
prs.save("example.pptx")
